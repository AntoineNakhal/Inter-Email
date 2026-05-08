"""PowerPoint (.pptx) text extraction via python-pptx.

We pull text from every shape that has a `text_frame` — title, body, tables,
and notes. Slide numbers are inserted so the chunker preserves locality
("Slide 4: ..." chunks together rather than running into Slide 5).
"""

from __future__ import annotations

import io

from backend.knowledge.extractors.base import BaseExtractor, ExtractionError


class PptxExtractor(BaseExtractor):
    file_type = "pptx"

    def extract(self, content: bytes, *, filename: str = "") -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ExtractionError(
                "python-pptx is not installed — run `pip install python-pptx`."
            ) from exc

        try:
            presentation = Presentation(io.BytesIO(content))
            slides_text: list[str] = []
            for slide_index, slide in enumerate(presentation.slides, start=1):
                fragments: list[str] = []
                for shape in slide.shapes:
                    text = self._shape_text(shape)
                    if text:
                        fragments.append(text)

                # Speaker notes — often where vendors hide setup details.
                notes_slide = getattr(slide, "notes_slide", None)
                if notes_slide and notes_slide.notes_text_frame:
                    note_text = notes_slide.notes_text_frame.text or ""
                    if note_text.strip():
                        fragments.append(f"[Notes] {note_text.strip()}")

                if fragments:
                    slides_text.append(
                        f"Slide {slide_index}:\n" + "\n".join(fragments)
                    )
        except Exception as exc:
            raise ExtractionError(
                f"Failed to extract text from PowerPoint '{filename}': {exc}"
            ) from exc

        return "\n\n".join(slides_text).strip()

    @staticmethod
    def _shape_text(shape) -> str:
        """Pull text from a shape, handling both text frames and tables."""
        # Tables: walk cells (each cell's text_frame).
        if getattr(shape, "has_table", False) and getattr(shape, "table", None):
            rows: list[str] = []
            for row in shape.table.rows:
                cells = [cell.text_frame.text.strip() for cell in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            return "\n".join(rows)

        # Plain text frames.
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is not None:
            return (text_frame.text or "").strip()
        return ""
